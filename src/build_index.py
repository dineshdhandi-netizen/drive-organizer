"""Build a SQLite FTS5 search index over text-bearing files in a folder tree.

Walks the target recursively, extracts text from PDF/DOCX/PPTX/XLSX/TXT/MD
files, and stores it in a SQLite database with FTS5 full-text search.

Incremental: re-running skips files whose path+mtime+size match an existing
row. Pass --rebuild to wipe and re-index.

Usage:
  python build_index.py "G:\\My Drive"
  python build_index.py "G:\\My Drive" --db reports/search.db --rebuild
"""

from __future__ import annotations

import argparse
import sqlite3
import sys
import time
from pathlib import Path


TEXT_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".rtf"}

# Folders to skip — dupes shouldn't be indexed twice; photo/video subfolders
# have no text. Extension filter handles the rest.
SKIP_DIR_NAMES = {"Duplicates", "Photos", "Videos", "Audio", "GoogleDocs"}


SCHEMA = """
CREATE TABLE IF NOT EXISTS documents (
    id        INTEGER PRIMARY KEY,
    path      TEXT UNIQUE NOT NULL,
    filename  TEXT NOT NULL,
    ext       TEXT NOT NULL,
    size      INTEGER NOT NULL,
    mtime     REAL NOT NULL,
    indexed_at REAL NOT NULL
);

CREATE VIRTUAL TABLE IF NOT EXISTS documents_fts USING fts5(
    filename, content,
    tokenize='unicode61'
);
"""


def extract_pdf(path: Path, page_cap: int = 30) -> str:
    from pypdf import PdfReader
    try:
        r = PdfReader(str(path))
        if r.is_encrypted:
            return ""
        out: list[str] = []
        for i, page in enumerate(r.pages):
            if i >= page_cap:
                break
            try:
                out.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(out)
    except Exception:
        return ""


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        return "\n".join(p.text for p in Document(str(path)).paragraphs)
    except Exception:
        return ""


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        out: list[str] = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out)
    except Exception:
        return ""


def extract_xlsx(path: Path, max_cells: int = 500) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        out: list[str] = []
        cells = 0
        for name in wb.sheetnames:
            out.append(name)
            ws = wb[name]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        out.append(str(cell))
                        cells += 1
                        if cells >= max_cells:
                            return "\n".join(out)
        return "\n".join(out)
    except Exception:
        return ""


def extract_plain(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":     return extract_pdf(path)
    if ext == ".docx":    return extract_docx(path)
    if ext == ".pptx":    return extract_pptx(path)
    if ext == ".xlsx":    return extract_xlsx(path)
    if ext in {".txt", ".md", ".rtf"}: return extract_plain(path)
    return ""


def iter_candidates(root: Path):
    for p in root.rglob("*"):
        if not p.is_file():
            continue
        # Skip if any path part is in SKIP_DIR_NAMES.
        if any(part in SKIP_DIR_NAMES for part in p.relative_to(root).parts[:-1]):
            continue
        if p.suffix.lower() not in TEXT_EXTS:
            continue
        yield p


def init_db(db_path: Path, rebuild: bool) -> sqlite3.Connection:
    db_path.parent.mkdir(parents=True, exist_ok=True)
    if rebuild and db_path.exists():
        db_path.unlink()
    conn = sqlite3.connect(str(db_path))
    conn.executescript(SCHEMA)
    return conn


def index_file(conn: sqlite3.Connection, path: Path) -> str:
    """Returns 'inserted', 'skipped', or 'error'."""
    try:
        st = path.stat()
    except OSError:
        return "error"
    cur = conn.execute(
        "SELECT id, size, mtime FROM documents WHERE path = ?",
        (str(path),),
    )
    existing = cur.fetchone()
    if existing and existing[1] == st.st_size and existing[2] == st.st_mtime:
        return "skipped"

    text = extract_text(path)
    if not text:
        # Still record the file so we know we tried.
        text = ""

    if existing:
        conn.execute("DELETE FROM documents WHERE id = ?", (existing[0],))
    conn.execute(
        """INSERT INTO documents (path, filename, ext, size, mtime, indexed_at)
           VALUES (?, ?, ?, ?, ?, ?)""",
        (str(path), path.name, path.suffix.lower(), st.st_size, st.st_mtime, time.time()),
    )
    row_id = conn.execute("SELECT last_insert_rowid()").fetchone()[0]
    conn.execute(
        "INSERT INTO documents_fts(rowid, filename, content) VALUES (?, ?, ?)",
        (row_id, path.name, text),
    )
    return "inserted"


def main() -> None:
    p = argparse.ArgumentParser(description="Build a SQLite FTS5 search index.")
    p.add_argument("target", type=Path, help="Root folder to index.")
    p.add_argument("--db", type=Path, default=Path("reports/search.db"))
    p.add_argument("--rebuild", action="store_true", help="Wipe and re-index everything.")
    p.add_argument("--limit", type=int, default=None, help="Index at most N files (testing).")
    args = p.parse_args()

    if not args.target.is_dir():
        sys.exit(f"Not a directory: {args.target}")

    conn = init_db(args.db, args.rebuild)
    inserted = skipped = errored = 0
    n = 0
    t0 = time.time()
    print(f"Indexing {args.target} -> {args.db}")
    for path in iter_candidates(args.target):
        n += 1
        try:
            result = index_file(conn, path)
        except Exception as e:
            print(f"  error: {path}  {e}", file=sys.stderr)
            errored += 1
            continue
        if result == "inserted":
            inserted += 1
        elif result == "skipped":
            skipped += 1
        if n % 50 == 0:
            elapsed = time.time() - t0
            rate = n / elapsed if elapsed > 0 else 0
            print(f"  {n} files seen | {inserted} new | {skipped} skipped | {rate:.1f}/s")
        if args.limit and n >= args.limit:
            break
        if n % 100 == 0:
            conn.commit()
    conn.commit()
    conn.close()
    elapsed = time.time() - t0
    print(f"\nDone in {elapsed:.1f}s.  scanned={n}  new={inserted}  skipped={skipped}  errored={errored}")
    print(f"Index: {args.db.resolve()}")


if __name__ == "__main__":
    main()
