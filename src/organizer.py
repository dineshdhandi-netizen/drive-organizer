"""Drive Organizer — sort loose top-level files into category folders.

Classification cascade (first match wins):
  1. Skip system files (desktop.ini, Thumbs.db, .DS_Store, dotfiles)
  2. Property — match by filename keywords
  3. Property — match by extracted file content (PDF/DOCX/XLSX/PPTX)
  4. Topic    — match by filename keywords
  5. Topic    — match by extracted file content
  6. Type fallback — Misc/<Type>/ (photos sub-folded by EXIF year)

Modes:
  (default)    dry-run: prints proposed moves, changes nothing
  --execute    actually move files
  --limit N    process only the first N candidate files (for safe test runs)
"""

from __future__ import annotations

import argparse
import csv
import re
import shutil
import sys
from collections import Counter
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path


# Destination paths are RELATIVE to the target root (e.g. G:\My Drive).
PROPERTY_RULES: list[tuple[str, tuple[str, ...], Path]] = [
    ("6525 Abrams Dr",      ("6525 abrams", "6525  abrams"),                       Path("6525 Abrams Dr")),
    ("2916 Puma",           ("2916 puma", "puma rd", "puma road"),                 Path("Auctus Realty LLC/2916 Puma")),
    ("2933 Noble Oaks Dr",  ("2933 noble", "noble oaks"),                          Path("Auctus Realty LLC/2933 Noble Oaks Dr")),
    ("3000 Duchess Trl",    ("3000 duchess", "duchess trl", "duchess trail"),      Path("Auctus Realty LLC/3000 Duchess Trl")),
    ("337 Metropolitan Dr", ("337 metropolitan", "metropolitan dr"),               Path("Auctus Realty LLC/337 Metropolitan Dr")),
    ("Hyderabad Flat #1",   ("hyderabad flat", "hyderabad apartment"),             Path("Auctus Realty LLC/Hyderabad Flat #1")),
]

# Phrases that LABEL a property reference. We only trust a content-based
# property match if a known property keyword appears within ~200 chars after
# one of these labels — this avoids false positives from mailing addresses.
PROPERTY_LABEL_PATTERNS: tuple[str, ...] = (
    r"property\s*address\b",
    r"covered\s*property\b",
    r"subject\s*property\b",
    r"property\s*location\b",
    r"address\s*of\s*property\b",
    r"insured\s*location\b",
    r"premises\s*address\b",
    r"property\s*owner\b",
)


# (topic, filename_keywords, content_keywords, destination)
# Order matters: first match wins. Family is before Legal so birth certs go to Family.
TOPIC_RULES: list[tuple[str, tuple[str, ...], tuple[str, ...], Path]] = [
    ("Taxes",          ("tax return", "w2", "w-2", "1099", "1040", " irs "),
                       ("form w-2", "form 1099", "internal revenue service", "tax year", "form 1040"),
                       Path("Tax Documents")),
    ("Resume",         ("resume", " cv ", "linkedin", "cover letter"),
                       ("curriculum vitae", "professional experience"),
                       Path("Resume")),
    ("AutoInsurance",  ("honda", "odyssey", "toyota", "tesla insurance", "auto insurance", "car insurance"),
                       ("vehicle identification number", "auto policy"),
                       Path("Car Insurance and Title")),
    ("BankOfAmerica",  (" boa ", "bank of america", "bofa"),
                       ("bank of america, n.a", "bank of america corporation"),
                       Path("Bank of America")),
    ("ICICI",          ("icici",),
                       ("icici bank",),
                       Path("Icici")),
    ("HSA",            (" hsa ", "health savings"),
                       ("health savings account",),
                       Path("HSA - Health Receipts")),
    ("Charity",        ("donation", "charity", "charitable", "501c3"),
                       ("tax-exempt", "501(c)(3)"),
                       Path("Charity Receipts")),
    ("Trading",        ("trading journal", "brokerage", "robinhood", "schwab", "fidelity", "etrade"),
                       ("brokerage account", "shares purchased"),
                       Path("Trading")),
    ("Health",         ("medical", "doctor", "prescription", "dental"),
                       ("patient name", "icd-10", "diagnosis code", "medical record"),
                       Path("Family Health")),
    ("Travel",         ("flight", "airline", "hotel", "trip", "visa", "passport", "itinerary"),
                       ("departure", "boarding pass", "confirmation number"),
                       Path("Travel")),
    ("Family",         ("birth certificate", "marriage certificate", "wedding", "birthday",
                        "b-day", "anniversary", "engagement", "stuti", "shloak", "baby"),
                       (),
                       Path("Family")),
    ("Legal",          ("uscis", "immigration", "court", "lawyer", "attorney", "certificate",
                        "drivers license", "driver's license", "driverlicense", "nda"),
                       ("plaintiff", "defendant", "non-disclosure agreement"),
                       Path("Legal")),
    ("Work",           ("qbr", "h2o.ai", "h2o.", "att-", "att+", "pay stub", "paystub",
                        "separation agreement", "offer letter", "client deck", "proposal"),
                       (),
                       Path("Work")),
    ("Banking",        ("bank statement", "account statement", "checking", "savings account"),
                       ("routing number",),
                       Path("Banking")),
    ("Insurance",      ("insurance",),
                       ("policy number", "named insured", "insurance carrier"),
                       Path("Insurance")),
]


TYPE_RULES: dict[str, tuple[str, ...]] = {
    "PDF":          (".pdf",),
    "Image":        (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp", ".heic"),
    "Document":     (".doc", ".docx", ".odt", ".rtf", ".txt", ".md"),
    "Spreadsheet":  (".xls", ".xlsx", ".csv", ".tsv", ".ods"),
    "Presentation": (".ppt", ".pptx", ".odp", ".key"),
    "Archive":      (".zip", ".rar", ".7z", ".tar", ".gz"),
    "Video":        (".mp4", ".mov", ".avi", ".mkv", ".webm"),
    "Audio":        (".mp3", ".wav", ".flac", ".aac", ".m4a"),
    "GoogleNative": (".gdoc", ".gsheet", ".gslides", ".gform", ".gdraw"),
}


TYPE_DEST: dict[str, Path] = {
    "PDF":          Path("Misc/PDFs"),
    "Document":     Path("Misc/Documents"),
    "Spreadsheet":  Path("Misc/Spreadsheets"),
    "Presentation": Path("Misc/Presentations"),
    "Archive":      Path("Misc/Archives"),
    "Video":        Path("Misc/Videos"),
    "Audio":        Path("Misc/Audio"),
    "GoogleNative": Path("Misc/GoogleDocs"),
    "Other":        Path("Misc/Other"),
    # Image handled specially (EXIF year).
}


SKIP_NAMES = {"desktop.ini", "thumbs.db", ".ds_store"}


def classify_type(extension: str) -> str:
    ext = extension.lower()
    for category, exts in TYPE_RULES.items():
        if ext in exts:
            return category
    return "Other"


def should_skip(name: str) -> bool:
    lower = name.lower()
    if lower in SKIP_NAMES:
        return True
    if name.startswith("."):
        return True
    return False


def match_property_by_filename(name: str) -> str | None:
    lower = name.lower()
    for prop, keywords, _ in PROPERTY_RULES:
        for kw in keywords:
            if kw in lower:
                return prop
    return None


def match_property_by_content(text: str) -> str | None:
    """Match a property ONLY when its keyword appears after a property-label phrase.

    Mailing addresses are everywhere on US documents, so a bare address mention
    is unreliable. Labels like 'Property Address:' or 'Covered Property:' are
    the authoritative signal.
    """
    lower = text.lower()
    for pattern in PROPERTY_LABEL_PATTERNS:
        for m in re.finditer(pattern, lower):
            window = lower[m.end(): m.end() + 200]
            for prop, keywords, _ in PROPERTY_RULES:
                for kw in keywords:
                    if kw in window:
                        return prop
            # Also try the 4 distinctive address numbers in the window
            # (so a content like 'Property Address: 2916 Puma Rd...' works even
            # if the keyword list doesn't include the bare number).
            for prop, keywords, _ in PROPERTY_RULES:
                for kw in keywords:
                    # Match just the leading number+street word (e.g. '2916 puma')
                    if any(token in window for token in kw.split() if len(token) >= 4):
                        # Require at least 2 tokens of the keyword in window to be safe.
                        tokens = kw.split()
                        if len(tokens) >= 2 and all(t in window for t in tokens):
                            return prop
    return None


def match_topic_in_text(text: str, by: str) -> str | None:
    """by = 'filename' uses filename_keywords; by = 'content' uses content_keywords."""
    lower = text.lower()
    for topic, fn_kws, ct_kws, _ in TOPIC_RULES:
        kws = fn_kws if by == "filename" else ct_kws
        for kw in kws:
            if kw in lower:
                return topic
    return None


def property_dest(name: str) -> Path:
    for prop, _, dest in PROPERTY_RULES:
        if prop == name:
            return dest
    raise KeyError(name)


def topic_dest(name: str) -> Path:
    for topic, _, _, dest in TOPIC_RULES:
        if topic == name:
            return dest
    raise KeyError(name)


# ---------- Content extractors (lazy imports) ----------

def extract_pdf(path: Path, page_cap: int = 20) -> str:
    from pypdf import PdfReader
    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            return ""
        chunks: list[str] = []
        for i, page in enumerate(reader.pages):
            if i >= page_cap:
                break
            try:
                chunks.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(chunks)
    except Exception:
        return ""


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks)
    except Exception:
        return ""


def extract_xlsx(path: Path, max_cells: int = 200) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        chunks: list[str] = []
        cells_read = 0
        for sheet_name in wb.sheetnames:
            chunks.append(sheet_name)
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        chunks.append(str(cell))
                        cells_read += 1
                        if cells_read >= max_cells:
                            return "\n".join(chunks)
        return "\n".join(chunks)
    except Exception:
        return ""


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".xlsx":
        return extract_xlsx(path)
    return ""


def exif_year(path: Path) -> int | None:
    try:
        from PIL import Image, ExifTags
    except Exception:
        return None
    try:
        with Image.open(path) as img:
            exif = img.getexif()
            if not exif:
                return None
            # 36867 = DateTimeOriginal; 306 = DateTime (modification)
            for tag_id in (36867, 306):
                value = exif.get(tag_id)
                if value:
                    m = re.match(r"(\d{4})", str(value))
                    if m:
                        year = int(m.group(1))
                        if 1990 <= year <= 2100:
                            return year
    except Exception:
        return None
    return None


def mtime_year(path: Path) -> int:
    return datetime.fromtimestamp(path.stat().st_mtime).year


# ---------- Classification ----------

@dataclass
class Move:
    source: Path
    dest_dir: Path     # absolute
    reason: str

    @property
    def dest_path(self) -> Path:
        return self.dest_dir / self.source.name


def classify(file: Path, root: Path) -> tuple[Path, str]:
    """Return (absolute_dest_dir, reason)."""
    name = file.name

    prop = match_property_by_filename(name)
    if prop:
        return (root / property_dest(prop), f"property-filename:{prop}")

    text = extract_text(file)
    if text:
        prop = match_property_by_content(text)
        if prop:
            return (root / property_dest(prop), f"property-content:{prop}")

    topic = match_topic_in_text(name, by="filename")
    if topic:
        return (root / topic_dest(topic), f"topic-filename:{topic}")

    if text:
        topic = match_topic_in_text(text, by="content")
        if topic:
            return (root / topic_dest(topic), f"topic-content:{topic}")

    category = classify_type(file.suffix)
    if category == "Image":
        year = exif_year(file) or mtime_year(file)
        return (root / Path("Misc/Photos") / str(year), f"type:Image (year={year})")
    return (root / TYPE_DEST[category], f"type:{category}")


# ---------- Move planning & execution ----------

def plan_moves(root: Path, limit: int | None) -> list[Move]:
    moves: list[Move] = []
    count = 0
    for entry in sorted(root.iterdir(), key=lambda p: p.name.lower()):
        if not entry.is_file():
            continue
        if should_skip(entry.name):
            continue
        dest_dir, reason = classify(entry, root)
        moves.append(Move(source=entry, dest_dir=dest_dir, reason=reason))
        count += 1
        if limit is not None and count >= limit:
            break
    return moves


def unique_dest(dest: Path) -> Path:
    """Append -1, -2, ... if destination already exists."""
    if not dest.exists():
        return dest
    stem, suffix = dest.stem, dest.suffix
    for i in range(1, 10000):
        candidate = dest.with_name(f"{stem}-{i}{suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"Too many collisions for {dest}")


def execute_moves(moves: list[Move], log_path: Path) -> tuple[int, int]:
    log_path.parent.mkdir(parents=True, exist_ok=True)
    succeeded = 0
    failed = 0
    is_new = not log_path.exists()
    with log_path.open("a", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        if is_new:
            writer.writerow(["timestamp", "filename", "source", "dest", "reason", "status"])
        now_iso = datetime.now(timezone.utc).isoformat(timespec="seconds")
        for mv in moves:
            try:
                mv.dest_dir.mkdir(parents=True, exist_ok=True)
                target = unique_dest(mv.dest_path)
                shutil.move(str(mv.source), str(target))
                writer.writerow([now_iso, mv.source.name, str(mv.source), str(target), mv.reason, "ok"])
                succeeded += 1
            except Exception as e:
                writer.writerow([now_iso, mv.source.name, str(mv.source), str(mv.dest_path), mv.reason, f"error: {e}"])
                failed += 1
    return succeeded, failed


def print_plan(moves: list[Move]) -> None:
    by_dest = Counter(str(mv.dest_dir) for mv in moves)
    by_reason_prefix = Counter(mv.reason.split(":")[0] for mv in moves)

    print(f"\nPlanned moves: {len(moves)}")
    print("\nBy destination (top 25):")
    for dest, n in by_dest.most_common(25):
        print(f"  {n:4}  {dest}")

    print("\nBy classifier:")
    for r, n in by_reason_prefix.most_common():
        print(f"  {n:4}  {r}")

    print("\nFirst 20 individual moves:")
    for mv in moves[:20]:
        print(f"  {mv.source.name}")
        print(f"      -> {mv.dest_dir}   [{mv.reason}]")


def main() -> None:
    p = argparse.ArgumentParser(description="Organize loose top-level files in a folder.")
    p.add_argument("target", type=Path, help="Folder whose top-level files to organize.")
    p.add_argument("--execute", action="store_true", help="Actually move files (otherwise dry-run).")
    p.add_argument("--limit", type=int, default=None, help="Process only first N candidate files.")
    p.add_argument("--log", type=Path, default=Path("reports/moves.csv"), help="Move-log CSV path.")
    p.add_argument("--plan-out", type=Path, default=Path("reports/plan.csv"), help="Where to write the full plan as CSV.")
    args = p.parse_args()

    if not args.target.is_dir():
        sys.exit(f"Not a directory: {args.target}")

    print(f"Target: {args.target}")
    print(f"Mode:   {'EXECUTE (real moves)' if args.execute else 'DRY-RUN (no changes)'}")
    if args.limit:
        print(f"Limit:  {args.limit} files")

    moves = plan_moves(args.target, args.limit)

    # Write full plan to CSV so user can review.
    args.plan_out.parent.mkdir(parents=True, exist_ok=True)
    with args.plan_out.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["filename", "source", "dest_dir", "reason"])
        for mv in moves:
            w.writerow([mv.source.name, str(mv.source), str(mv.dest_dir), mv.reason])
    print(f"\nFull plan written to: {args.plan_out.resolve()}")

    print_plan(moves)

    if args.execute:
        print("\n--- EXECUTING MOVES ---")
        ok, fail = execute_moves(moves, args.log)
        print(f"Done. Succeeded: {ok}   Failed: {fail}")
        print(f"Log appended to: {args.log.resolve()}")
    else:
        print("\n(Dry run — no files moved.)")


if __name__ == "__main__":
    main()
